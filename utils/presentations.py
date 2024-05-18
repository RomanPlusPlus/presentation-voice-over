from pptx import Presentation
from pdf2image import convert_from_path
import os


def extract_notes_from_pptx(pptx_file):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    # List to store extracted notes
    notes = []

    # Iterate through each slide
    for slide in presentation.slides:
        # Initialize a string to store notes for each slide
        slide_notes = ""

        # Get the notes slide
        notes_slide = slide.notes_slide

        # Extract notes from the notes slide
        for paragraph in notes_slide.notes_text_frame.paragraphs:
            # Extract text from each paragraph
            slide_notes += paragraph.text + "\n"

        # Append slide notes to the notes list
        notes.append(slide_notes)

    return notes


def pdf_to_images(pdf_path, output_folder, dpi=300):
    """
    Convert a PDF into a series of images.

    :param pdf_path: Path to the input PDF file.
    :param output_folder: Folder to save the output images.
    :param dpi: Resolution of the output images.
    """
    print(f"Converting {pdf_path} to images...")
    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Convert PDF to images
    images = convert_from_path(pdf_path, dpi=dpi)

    # Save each image
    for i, image in enumerate(images):
        image_path = os.path.join(output_folder, f"slide_{i}.png")
        image.save(image_path, "PNG")
        print(f"Saved: {image_path}")
